"""PPTX extractor backed by ``python-pptx``.

Lazily imports the ``pptx`` package on first use; the import error message
points the user at the ``[ingest-pptx]`` extra. One section per slide; the
slide title placeholder becomes the heading, remaining text frames the body.
"""

from __future__ import annotations

import io
from typing import Any, ClassVar

from ...exceptions import KaguraIngestError
from .._types import ExtractedContent, ExtractedSection
from ._util import MAX_TOTAL_TEXT_CHARS as _MAX_TOTAL_TEXT_CHARS
from ._util import filename_title

# PPTX is a ZIP container — cap slides and total serialized text to defuse
# decompression bombs.
_MAX_SLIDES = 50_000


def _load_pptx() -> Any:
    try:
        import pptx  # type: ignore[import-untyped]
    except ImportError as e:
        raise KaguraIngestError(
            "python-pptx is not installed. Install with: pip install 'kagura-memory[ingest-pptx]'"
        ) from e
    return pptx


class PptxExtractor:
    """Structural PPTX extractor.

    Each slide becomes one section. The slide's title placeholder is the
    heading (falling back to ``Slide N``); the remaining shape text frames
    form the body. ``page_range`` is ``None`` — slides map to sections, not
    pages.
    """

    supports: ClassVar[frozenset[str]] = frozenset(
        {"application/vnd.openxmlformats-officedocument.presentationml.presentation"}
    )

    def extract(self, source: bytes, source_uri: str | None = None) -> ExtractedContent:
        pptx = _load_pptx()
        try:
            presentation = pptx.Presentation(io.BytesIO(source))
        except Exception as e:  # noqa: BLE001 - re-raised as our domain error
            raise KaguraIngestError(f"failed to open PPTX: {e}") from e

        # Check the slide count BEFORE materializing the list — python-pptx's
        # slide collection is sized (counts <sldId> elements) without building
        # Slide objects, so a hostile deck is rejected before allocation.
        slide_count = len(presentation.slides)
        if slide_count > _MAX_SLIDES:
            raise KaguraIngestError(f"PPTX slide count {slide_count} exceeds limit {_MAX_SLIDES}")
        slides = list(presentation.slides)

        sections = self._sections(slides)
        # Use the first slide's *real* title placeholder for the document
        # title — not the synthetic "Slide N" heading, which would shadow the
        # filename fallback for title-less decks.
        real_first_title = self._slide_text(slides[0])[0] if slides else None
        title = real_first_title or filename_title(source_uri)
        return ExtractedContent(title=title, sections=sections, images=[], page_count=None)

    @classmethod
    def _sections(cls, slides: list[Any]) -> list[ExtractedSection]:
        sections: list[ExtractedSection] = []
        total_chars = 0
        for idx, slide in enumerate(slides, start=1):
            heading, body = cls._slide_text(slide)
            heading = heading or f"Slide {idx}"
            total_chars += len(heading) + len(body)
            if total_chars > _MAX_TOTAL_TEXT_CHARS:
                raise KaguraIngestError(
                    f"PPTX total text exceeds {_MAX_TOTAL_TEXT_CHARS} chars (decompression bomb?)"
                )
            sections.append(
                ExtractedSection(
                    heading=heading,
                    body_text=body,
                    page_range=None,
                    depth=1,
                    anchor=heading,
                )
            )
        return sections

    @staticmethod
    def _slide_text(slide: Any) -> tuple[str | None, str]:
        title_text: str | None = None
        title_shape = None
        try:
            title_shape = slide.shapes.title
        except Exception:  # noqa: BLE001 - some layouts have no title placeholder
            title_shape = None
        if title_shape is not None and getattr(title_shape, "has_text_frame", False):
            title_text = (title_shape.text or "").strip() or None

        body_parts: list[str] = []
        for shape in slide.shapes:
            if shape is title_shape:
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            text = (shape.text or "").strip()
            if text:
                body_parts.append(text)
        return title_text, "\n".join(body_parts).strip()
